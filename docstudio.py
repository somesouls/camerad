# -*- coding: utf-8 -*-
"""Studio Dokumen (Epik C) — parser dokumen + chunking + generator output.

Modul MANDIRI: tidak mengimpor FastAPI maupun llm_client. web_app.py yang
mengorkestrasi pemanggilan LLM. Semua keluaran = dokumen untuk diunduh/preview
(bukan tulis balik ke DB). Jawaban dibatasi HANYA dari dokumen terunggah + DB
global (konsisten dengan Epik B).
"""
import io
import os
import re
import csv
import json

# Ekstensi yang didukung (minimal sesuai rencana Epik C)
IMAGE_EXTS = ["png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp", "gif"]
SUPPORTED_EXTS = ["txt", "md", "csv", "xlsx", "xls", "pdf", "docx", "pptx"] + IMAGE_EXTS

OUTPUT_TYPES = ["ringkasan", "laporan", "mindmap", "tabel"]

# Batas ukuran teks agar aman untuk konteks LLM.
MAX_CHUNK_CHARS = 6000
MAX_TABLE_ROWS = 5000
PREVIEW_CHARS = 1500


# =============================================================
# Guardrail & prompt builders (isi jawaban HANYA dari dokumen + DB global)
# =============================================================
GUARDRAIL = (
    "Kamu asisten Studio Dokumen internal Camerad Studio. Tugasmu mengolah ISI "
    "DOKUMEN yang diunggah pengguna. Jawab HANYA berdasarkan KONTEKS DOKUMEN "
    "(dan konteks internal bila disediakan) di bawah. DILARANG memakai "
    "pengetahuan umum/eksternal, mencari di web, atau mengarang fakta yang tidak "
    "ada di dokumen. Jika informasi yang diminta tidak ada di dokumen, katakan "
    "jujur bahwa dokumen tidak memuat informasi itu. Tulis dalam Bahasa "
    "Indonesia yang jelas dan rapi (boleh Markdown)."
)

_OUT_INSTRUCTION = {
    "ringkasan": (
        "Buat RINGKASAN yang padat dan setia pada isi dokumen. Susun: satu "
        "paragraf ikhtisar, lalu poin-poin kunci sebagai bullet. Jangan menambah "
        "informasi di luar dokumen."
    ),
    "laporan": (
        "Buat LAPORAN terstruktur berbasis isi dokumen dengan bagian: \n"
        "## Ringkasan Eksekutif\n## Temuan Utama\n## Rincian/Analisis\n"
        "## Kesimpulan & Rekomendasi. Gunakan bullet & tabel Markdown bila perlu. "
        "Semua isi harus bersumber dari dokumen; jangan mengarang angka."
    ),
    "mindmap": (
        "Buat KERANGKA (outline) hierarkis dari isi dokumen sebagai daftar "
        "bertingkat memakai tanda '-' dengan indentasi 2 spasi per tingkat. "
        "Maksimal 4 tingkat. Baris pertama = judul pusat (satu '-' tanpa "
        "indentasi). Jangan tulis penjelasan lain, HANYA outline."
    ),
    "tabel": (
        "Ekstrak data dari dokumen menjadi SATU tabel. Keluarkan HANYA JSON valid "
        'berbentuk {"columns": [...], "rows": [[...], ...]} tanpa teks lain, tanpa '
        "blok kode. Nilai sel berupa string/angka. Jika dokumen tidak memuat data "
        'tabular, keluarkan {"columns": [], "rows": []}.'
    ),
}


def map_prompt(chunk_text, question=""):
    """Prompt tahap MAP: ekstrak fakta relevan dari satu potongan dokumen."""
    q = ("\nFokus pada pertanyaan pengguna: " + question) if question else ""
    return (
        "Berikut satu bagian dari dokumen. Ekstrak poin-poin fakta penting apa "
        "adanya (ringkas, tanpa opini, tanpa menambah info luar)." + q +
        "\n\n=== BAGIAN DOKUMEN ===\n" + chunk_text
    )


def reduce_prompt(output_type, partials, question="", global_ctx="", filename=""):
    """Prompt tahap REDUCE: gabung hasil map menjadi keluaran akhir."""
    instr = _OUT_INSTRUCTION.get(output_type, _OUT_INSTRUCTION["ringkasan"])
    parts = []
    parts.append("Nama dokumen: " + (filename or "(tanpa nama)"))
    if question:
        parts.append("Permintaan/pertanyaan pengguna: " + question)
    parts.append("Instruksi keluaran: " + instr)
    if global_ctx:
        parts.append(
            "\n=== KONTEKS INTERNAL (pustaka/DB global, hanya untuk konsistensi "
            "istilah; jangan jadikan sumber data utama) ===\n" + global_ctx
        )
    body = "\n\n".join("[Bagian %d]\n%s" % (i + 1, p) for i, p in enumerate(partials))
    parts.append("\n=== RINGKASAN ISI DOKUMEN (hasil ekstraksi bertahap) ===\n" + body)
    return "\n\n".join(parts)


def single_prompt(output_type, doc_text, question="", global_ctx="", filename=""):
    """Prompt untuk dokumen kecil (tanpa map-reduce)."""
    instr = _OUT_INSTRUCTION.get(output_type, _OUT_INSTRUCTION["ringkasan"])
    parts = ["Nama dokumen: " + (filename or "(tanpa nama)")]
    if question:
        parts.append("Permintaan/pertanyaan pengguna: " + question)
    parts.append("Instruksi keluaran: " + instr)
    if global_ctx:
        parts.append(
            "\n=== KONTEKS INTERNAL (pustaka/DB global, untuk konsistensi istilah) "
            "===\n" + global_ctx
        )
    parts.append("\n=== ISI DOKUMEN ===\n" + doc_text)
    return "\n\n".join(parts)


# =============================================================
# Chunking
# =============================================================
def chunk_text(text, max_chars=MAX_CHUNK_CHARS):
    text = text or ""
    if len(text) <= max_chars:
        return [text] if text.strip() else []
    paras = re.split(r"\n\s*\n", text)
    chunks, cur = [], ""
    for p in paras:
        if len(p) > max_chars:
            # paragraf raksasa: pecah keras
            if cur:
                chunks.append(cur); cur = ""
            for i in range(0, len(p), max_chars):
                chunks.append(p[i:i + max_chars])
            continue
        if len(cur) + len(p) + 2 > max_chars:
            if cur:
                chunks.append(cur)
            cur = p
        else:
            cur = (cur + "\n\n" + p) if cur else p
    if cur.strip():
        chunks.append(cur)
    return chunks


# =============================================================
# Parser dokumen -> teks + tabel
# =============================================================
def _ext_of(filename):
    return os.path.splitext(filename or "")[1].lstrip(".").lower()


def _clean(s):
    return re.sub(r"[ \t]+", " ", (s or "").replace("\r", "")).strip()


def _rows_to_table(name, aoa):
    aoa = [[("" if c is None else str(c)) for c in row] for row in aoa if row is not None]
    aoa = [r for r in aoa if any((c or "").strip() for c in r)]
    if not aoa:
        return None
    cols = aoa[0]
    rows = aoa[1:MAX_TABLE_ROWS + 1]
    return {"name": name or "Tabel", "columns": cols, "rows": rows}


def _table_to_text(tbl, max_rows=40):
    lines = ["[Tabel: %s]" % tbl.get("name", "")]
    cols = tbl.get("columns") or []
    if cols:
        lines.append(" | ".join(str(c) for c in cols))
    for r in (tbl.get("rows") or [])[:max_rows]:
        lines.append(" | ".join(str(c) for c in r))
    extra = len(tbl.get("rows") or []) - max_rows
    if extra > 0:
        lines.append("... (%d baris lagi)" % extra)
    return "\n".join(lines)


def parse_txt(data):
    return data.decode("utf-8", errors="replace"), []


def parse_csv(data):
    txt = data.decode("utf-8", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(txt[:4096], delimiters=",;\t|")
    except Exception:
        dialect = csv.excel
    reader = csv.reader(io.StringIO(txt), dialect)
    aoa = [row for row in reader]
    tbl = _rows_to_table("CSV", aoa)
    tables = [tbl] if tbl else []
    text = _table_to_text(tbl) if tbl else txt
    return text, tables


def parse_xlsx(data):
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    tables, texts = [], []
    for ws in wb.worksheets:
        aoa = []
        for row in ws.iter_rows(values_only=True):
            aoa.append(list(row))
            if len(aoa) > MAX_TABLE_ROWS + 5:
                break
        tbl = _rows_to_table(ws.title, aoa)
        if tbl:
            tables.append(tbl)
            texts.append(_table_to_text(tbl))
    try:
        wb.close()
    except Exception:
        pass
    return "\n\n".join(texts), tables


# =============================================================
# OCR (opsional) — aktif bila Tesseract OCR terpasang di server
# =============================================================
def _ocr_lang():
    """Pilih bahasa OCR yang tersedia (utamakan ind+eng)."""
    try:
        import pytesseract
        langs = set(pytesseract.get_languages(config=""))
        picks = [l for l in ("ind", "eng") if l in langs]
        return "+".join(picks) if picks else None
    except Exception:
        return None


def ocr_available():
    """True bila pustaka OCR + biner Tesseract siap dipakai."""
    try:
        import pytesseract  # noqa
        from PIL import Image  # noqa
        import fitz  # noqa
    except Exception:
        return False
    try:
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def _ocr_image_obj(img):
    import pytesseract
    lang = _ocr_lang()
    try:
        return pytesseract.image_to_string(img, lang=lang) if lang else pytesseract.image_to_string(img)
    except Exception:
        return pytesseract.image_to_string(img)


def ocr_pdf(data, max_pages=40, zoom=2.0):
    """Render tiap halaman PDF ke gambar lalu OCR. Butuh Tesseract terpasang."""
    import fitz  # PyMuPDF
    from PIL import Image
    doc = fitz.open(stream=data, filetype="pdf")
    mat = fitz.Matrix(zoom, zoom)
    parts = []
    for i, pg in enumerate(doc):
        if i >= max_pages:
            break
        pix = pg.get_pixmap(matrix=mat)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        parts.append(_ocr_image_obj(img))
    n = doc.page_count
    doc.close()
    return "\n\n".join(parts), n


def ocr_image(data):
    """OCR satu berkas gambar. Butuh Tesseract terpasang."""
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    return _ocr_image_obj(img)


def _pdf_text_layer(data):
    """Ekstrak lapisan teks PDF via pypdf & PyMuPDF; ambil yang lebih lengkap."""
    t1, p1 = "", 0
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        t1 = "\n\n".join((pg.extract_text() or "") for pg in reader.pages)
        p1 = len(reader.pages)
    except Exception:
        pass
    t2, p2 = "", 0
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=data, filetype="pdf")
        t2 = "\n\n".join(pg.get_text("text") for pg in doc)
        p2 = doc.page_count
        doc.close()
    except Exception:
        pass
    text = t1 if len(t1.strip()) >= len(t2.strip()) else t2
    return text, max(p1, p2)


def parse_pdf(data):
    """Return (text, tables, pages, ocr_used).

    Utamakan lapisan teks (pypdf + PyMuPDF, ambil yang lebih lengkap). Bila
    hampir tanpa teks (kemungkinan hasil pindai) DAN OCR tersedia, jalankan OCR.
    """
    text, pages = _pdf_text_layer(data)
    ocr_used = False
    # Ambang: teks per-halaman sangat sedikit => kemungkinan PDF pindai.
    thin = len(text.strip()) < max(40, (pages or 1) * 8)
    if thin and ocr_available():
        try:
            ot, on = ocr_pdf(data)
            if len(ot.strip()) > len(text.strip()):
                text, pages, ocr_used = ot, (pages or on), True
        except Exception:
            pass
    return text, [], pages, ocr_used


def parse_docx(data):
    import docx  # python-docx
    d = docx.Document(io.BytesIO(data))
    paras = [p.text for p in d.paragraphs if (p.text or "").strip()]
    tables = []
    for ti, t in enumerate(d.tables):
        aoa = [[c.text for c in row.cells] for row in t.rows]
        tbl = _rows_to_table("Tabel %d" % (ti + 1), aoa)
        if tbl:
            tables.append(tbl)
    text = "\n\n".join(paras)
    if tables:
        text += "\n\n" + "\n\n".join(_table_to_text(t) for t in tables)
    return text, tables


def parse_pptx(data):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        buf = ["[Slide %d]" % (i + 1)]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs) or para.text
                    if (line or "").strip():
                        buf.append(line)
        slides.append("\n".join(buf))
    return "\n\n".join(slides), [], len(prs.slides)


def extract(data, filename):
    """Ekstrak teks + tabel dari file. Return dict.
    {ok, ext, text, tables, n_chars, pages, note}
    """
    ext = _ext_of(filename)
    if ext not in SUPPORTED_EXTS:
        return {"ok": False, "error": "Format .%s belum didukung. Didukung: %s" %
                (ext, ", ".join(SUPPORTED_EXTS))}
    pages = 0
    note = ""
    try:
        if ext in ("txt", "md"):
            text, tables = parse_txt(data)
        elif ext == "csv":
            text, tables = parse_csv(data)
        elif ext in ("xlsx", "xls"):
            text, tables = parse_xlsx(data)
        elif ext == "pdf":
            text, tables, pages, ocr_used = parse_pdf(data)
            if len(text.strip()) < 20:
                if ocr_available():
                    note = ("PDF hampir tanpa teks dan OCR tidak menghasilkan teks terbaca "
                            "— dokumen mungkin kosong atau kualitas pindai terlalu rendah.")
                else:
                    note = ("PDF nyaris tanpa teks — kemungkinan hasil pindai (gambar). "
                            "OCR belum aktif di server: pasang Tesseract OCR (lihat CHANGES_v11) "
                            "lalu unggah ulang.")
            elif ocr_used:
                note = "Teks diekstrak via OCR (dokumen terdeteksi hasil pindai/gambar)."
        elif ext in IMAGE_EXTS:
            if not ocr_available():
                return {"ok": False, "error": ("Berkas gambar membutuhkan OCR. Pasang Tesseract OCR "
                        "di server (lihat CHANGES_v11), lalu unggah ulang.")}
            text = ocr_image(data)
            tables = []
            pages = 1
            if len(text.strip()) < 3:
                note = "OCR tidak menemukan teks terbaca pada gambar ini."
            else:
                note = "Teks diekstrak dari gambar via OCR."
        elif ext == "docx":
            text, tables = parse_docx(data)
        elif ext == "pptx":
            text, tables, pages = parse_pptx(data)
        else:
            return {"ok": False, "error": "Format .%s belum didukung." % ext}
    except ImportError as e:
        return {"ok": False, "error": "Pustaka parser untuk .%s belum terpasang: %s" % (ext, e)}
    except Exception as e:
        return {"ok": False, "error": "Gagal membaca .%s: %s" % (ext, e)}

    text = (text or "").strip()
    return {
        "ok": True,
        "ext": ext,
        "text": text,
        "tables": tables or [],
        "n_chars": len(text),
        "pages": pages,
        "note": note,
    }


# =============================================================
# Konversi keluaran -> dokumen unduhan
# =============================================================
def outline_to_mermaid(outline):
    """Ubah outline bertingkat ('-' + indentasi 2 spasi) menjadi mermaid mindmap."""
    lines = [ln for ln in (outline or "").splitlines() if ln.strip()]
    out = ["mindmap"]
    root_done = False
    for ln in lines:
        m = re.match(r"^(\s*)[-*]\s+(.*)$", ln)
        if not m:
            m2 = re.match(r"^(\s*)(.*)$", ln)
            indent = len(m2.group(1)); label = m2.group(2).strip()
        else:
            indent = len(m.group(1)); label = m.group(2).strip()
        if not label:
            continue
        label = re.sub(r"[\(\)\[\]\{\}]", " ", label).strip()
        level = indent // 2
        if not root_done:
            out.append("  root((%s))" % label)
            root_done = True
        else:
            out.append("  " * (level + 2) + label)
    return "\n".join(out) if len(out) > 1 else "mindmap\n  root((Dokumen))"


def parse_table_json(s):
    """Parse keluaran LLM menjadi (columns, rows) dengan toleransi blok kode."""
    if not s:
        return [], []
    t = s.strip()
    t = re.sub(r"^```[a-zA-Z]*\n?", "", t)
    t = re.sub(r"\n?```$", "", t).strip()
    m = re.search(r"\{.*\}", t, re.S)
    if m:
        t = m.group(0)
    try:
        d = json.loads(t)
    except Exception:
        return [], []
    cols = d.get("columns") or []
    rows = d.get("rows") or []
    cols = [str(c) for c in cols]
    norm = []
    for r in rows:
        if isinstance(r, dict):
            norm.append([str(r.get(c, "")) for c in cols])
        elif isinstance(r, (list, tuple)):
            norm.append([("" if c is None else str(c)) for c in r])
    return cols, norm


def table_to_csv_bytes(columns, rows):
    buf = io.StringIO()
    w = csv.writer(buf)
    if columns:
        w.writerow(columns)
    for r in rows:
        w.writerow(r)
    return buf.getvalue().encode("utf-8-sig")


def table_to_xlsx_bytes(columns, rows, sheet="Data"):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = (sheet or "Data")[:31]
    if columns:
        ws.append([str(c) for c in columns])
    for r in rows:
        ws.append([("" if c is None else str(c)) for c in r])
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def table_to_html(columns, rows, limit=200):
    def esc(x):
        return (str(x) if x is not None else "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    out = ["<table><thead><tr>"]
    for c in columns:
        out.append("<th>%s</th>" % esc(c))
    out.append("</tr></thead><tbody>")
    for r in rows[:limit]:
        out.append("<tr>" + "".join("<td>%s</td>" % esc(v) for v in r) + "</tr>")
    out.append("</tbody></table>")
    return "".join(out)


def md_to_bytes(text):
    return (text or "").encode("utf-8")


def pick_largest_table(tables):
    best, best_cells = None, -1
    for t in (tables or []):
        cells = len(t.get("columns") or []) * (len(t.get("rows") or []) + 1)
        if cells > best_cells:
            best, best_cells = t, cells
    return best

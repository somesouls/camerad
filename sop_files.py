# -*- coding: utf-8 -*-
"""sop_files.py — Ekstraksi teks dokumen SOP/Proses Bisnis + pemotongan per-bagian.

Menerima berbagai format dokumen dan mengembalikan:
  * teks lengkap,
  * daftar BAGIAN (section) hasil pemotongan berbasis heading -> tiap bagian
    kelak menjadi satu unit sop_unit (chunk untuk RAG),
  * kandidat JUDUL dokumen (dari metadata/heading/nama berkas),
  * kandidat KATEGORI (SOP / Proses Bisnis / Panduan) dari kata kunci.

Format yang didukung:
  * PDF          -> teks (PyMuPDF) atau OCR bila scan (do_ocr) — via peraturan_files
  * PPTX         -> python-pptx: tiap slide = satu bagian (judul slide = heading)
  * DOCX         -> python-docx: pecah pada paragraf bergaya Heading/Judul + tabel
  * TXT / MD     -> teks polos, dipotong via heuristik heading
  * HTML / HTM   -> BeautifulSoup: pecah pada tag h1..h4
  * Gambar       -> OCR (jpg/png/tif/...) bila do_ocr

OCR (PDF scan & gambar) DIPAKAI ULANG dari peraturan_files (butuh biner
Tesseract 'ind' + Poppler; bila tak ada, berkas ditandai perlu_ocr). Progres
OCR memakai mekanisme callback peraturan_files.set_progress_cb yang sama, jadi
UI SOP bisa memantau OCR seperti menu Peraturan.

Format lama .ppt/.doc (biner OLE) TIDAK didukung python-pptx/docx -> ditandai
agar dikonversi manual ke .pptx/.docx.
"""
from __future__ import annotations

import os
import re
from dataclasses import dataclass, field

import peraturan_files as PF  # dipakai ulang: OCR PDF/gambar + progress cb

try:
    from bs4 import BeautifulSoup
except Exception:  # pragma: no cover
    BeautifulSoup = None  # type: ignore

# Ekstensi yang didukung untuk impor SOP.
DOC_EXT = (".pdf", ".pptx", ".docx", ".txt", ".md", ".html", ".htm")
IMG_EXT = getattr(PF, "IMG_EXT", (".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp", ".gif"))
SKIP_NAMES = getattr(PF, "SKIP_NAMES", {"thumbs.db", "desktop.ini", ".ds_store"})
# Format lama biner yang perlu dikonversi dulu.
LAWAS_EXT = (".ppt", ".doc")

# Panjang maksimum satu bagian sebelum dipotong lagi (jaga chunk RAG tetap fokus).
MAKS_ISI = 1800
MIN_ISI = 25


@dataclass
class SopFileInfo:
    path: str
    tipe: str = "unknown"        # pdf | pptx | docx | txt | html | image | lawas | unknown
    judul_kandidat: str = ""
    kategori_kandidat: str = ""
    sections: list = field(default_factory=list)   # [{'bagian': str, 'isi': str}]
    teks: str = ""
    perlu_ocr: bool = False
    catatan: str = ""


def _clean(t: str) -> str:
    if not t:
        return ""
    t = t.replace("\xa0", " ")
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def _n_nonspace(t: str) -> int:
    return len(re.sub(r"\s+", "", t or ""))


# ---------------------------------------------------------------- judul & kategori
_HASHY = re.compile(r"^[0-9a-f]{8,}$", re.I)


def _nama_dari_file(path: str) -> str:
    base = os.path.splitext(os.path.basename(path))[0]
    base = re.sub(r"[_\-]+", " ", base)
    base = re.sub(r"\s+", " ", base).strip()
    return base


def _bersih_judul(s: str) -> str:
    s = _clean(s or "")
    s = s.strip(" -:•·–—\t")
    return s[:180]


def judul_bermakna(s: str) -> bool:
    """True bila kandidat judul layak dipakai apa adanya (tanpa bantuan AI)."""
    s = (s or "").strip()
    if len(s) < 6:
        return False
    if _HASHY.match(s.replace(" ", "")):
        return False
    huruf = sum(1 for c in s if c.isalpha())
    kata = len(s.split())
    return huruf >= 4 and kata >= 2


def tebak_kategori(nama: str, teks: str) -> str:
    hay = ((nama or "") + " " + (teks or "")[:2500]).lower()
    if any(x in hay for x in ("proses bisnis", "business process", "alur kerja", "workflow", "bagan alir", "flowchart")):
        return "Proses Bisnis"
    if "standar operasional" in hay or re.search(r"\bsop\b", hay):
        return "SOP"
    if any(x in hay for x in ("panduan", "petunjuk teknis", "juknis", "pedoman", "manual", "tata cara")):
        return "Panduan"
    return ""


# ---------------------------------------------------------------- heuristik heading
_RE_NUM = re.compile(r"^\s*\d+(\.\d+)*[\.\)]?\s+\S")
_RE_HURUF = re.compile(r"^\s*[A-Za-z][\.\)]\s+\S")
_RE_KATA_HEAD = re.compile(
    r"^\s*(BAB|BAGIAN|LANGKAH|TAHAP(?:AN)?|PROSEDUR|PENDAHULUAN|TUJUAN|RUANG\s+LINGKUP|"
    r"DEFINISI|REFERENSI|KEBIJAKAN|PROSES|ALUR|DASAR\s+HUKUM|LAMPIRAN)\b",
    re.I,
)


def _is_heading(line: str) -> bool:
    s = (line or "").strip()
    if not s or len(s) > 90:
        return False
    if _RE_KATA_HEAD.match(s):
        return True
    if _RE_NUM.match(s) and len(s) <= 90:
        return True
    if _RE_HURUF.match(s) and len(s) <= 80:
        return True
    if s.endswith(":") and len(s) <= 70 and len(s.split()) <= 9:
        return True
    # baris pendek HURUF BESAR semua (judul bab)
    huruf = [c for c in s if c.isalpha()]
    if huruf and len(s) <= 70 and s.upper() == s and len(s.split()) <= 10:
        return True
    return False


def seksi_dari_teks(teks: str) -> list:
    """Pecah teks polos menjadi bagian berbasis heuristik heading."""
    lines = (teks or "").splitlines()
    sections = []
    cur_head = ""
    cur = []

    def flush():
        isi = _clean("\n".join(cur))
        if _n_nonspace(isi) >= MIN_ISI or (isi and cur_head):
            sections.append({"bagian": cur_head, "isi": isi})

    for ln in lines:
        s = ln.strip()
        if s and _is_heading(s):
            flush()
            cur = []
            cur_head = s.rstrip(":").strip()
        else:
            cur.append(ln)
    flush()
    if not sections:
        isi = _clean(teks)
        if isi:
            sections = [{"bagian": "", "isi": isi}]
    return sections


# ------------------------------------------------------------------- ekstraktor
def _extract_pdf(path, do_ocr):
    teks, n, perlu = PF.extract_pdf(path, do_ocr=do_ocr)
    return teks or "", perlu


def _extract_image(path, do_ocr):
    teks, perlu = PF.extract_image(path, do_ocr=do_ocr)
    return teks or "", perlu


def _extract_txt(path):
    try:
        return open(path, encoding="utf-8", errors="replace").read()
    except Exception:
        return ""


def _extract_pptx(path):
    """(judul, sections, teks_gabungan). Tiap slide -> satu bagian."""
    try:
        from pptx import Presentation
    except Exception:
        return "", [], "", "python-pptx belum terpasang"
    try:
        prs = Presentation(path)
    except Exception as e:
        return "", [], "", "gagal buka pptx: %s" % str(e)[:120]
    sections, all_txt = [], []
    judul = ""
    for idx, slide in enumerate(prs.slides, start=1):
        head = ""
        try:
            if slide.shapes.title is not None and (slide.shapes.title.text or "").strip():
                head = slide.shapes.title.text.strip()
        except Exception:
            head = ""
        body = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if not t and para.text:
                            t = para.text.strip()
                        if t and t != head:
                            body.append(t)
                elif shape.has_table:
                    for row in shape.table.rows:
                        cells = [(_c.text or "").strip() for _c in row.cells]
                        line = " | ".join(c for c in cells if c)
                        if line:
                            body.append(line)
            except Exception:
                continue
        if idx == 1 and head:
            judul = head
        isi = _clean("\n".join(body))
        bagian = head or ("Slide %d" % idx)
        if _n_nonspace(isi) >= MIN_ISI or head:
            sections.append({"bagian": bagian, "isi": isi})
            all_txt.append((head + "\n" + isi).strip())
    return judul, sections, _clean("\n\n".join(all_txt)), ""


def _extract_docx(path):
    """(judul, sections, teks_gabungan). Pecah pada paragraf gaya Heading."""
    try:
        from docx import Document
    except Exception:
        return "", [], "", "python-docx belum terpasang"
    try:
        doc = Document(path)
    except Exception as e:
        return "", [], "", "gagal buka docx: %s" % str(e)[:120]
    judul = ""
    try:
        judul = (doc.core_properties.title or "").strip()
    except Exception:
        judul = ""
    sections, all_txt = [], []
    cur_head, cur = "", []

    def flush():
        isi = _clean("\n".join(cur))
        if _n_nonspace(isi) >= MIN_ISI or (isi and cur_head):
            sections.append({"bagian": cur_head, "isi": isi})
            all_txt.append((cur_head + "\n" + isi).strip())

    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        style = ""
        try:
            style = (p.style.name if p.style else "") or ""
        except Exception:
            style = ""
        sl = style.lower()
        is_head = sl.startswith("heading") or sl.startswith("judul") or sl.startswith("title")
        if is_head and txt:
            flush()
            cur = []
            cur_head = txt
            if not judul:
                judul = txt
        elif txt:
            cur.append(txt)
    flush()
    # Tabel (mis. matriks RACI / langkah proses)
    try:
        for ti, tbl in enumerate(doc.tables, start=1):
            rows = []
            for row in tbl.rows:
                cells = [(_c.text or "").strip() for _c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    rows.append(line)
            isi = _clean("\n".join(rows))
            if _n_nonspace(isi) >= MIN_ISI:
                sections.append({"bagian": "Tabel %d" % ti, "isi": isi})
                all_txt.append(isi)
    except Exception:
        pass
    return judul, sections, _clean("\n\n".join(all_txt)), ""


def _extract_html(path):
    """(judul, sections, teks_gabungan). Pecah pada tag h1..h4."""
    try:
        html = open(path, encoding="utf-8", errors="replace").read()
    except Exception as e:
        return "", [], "", "gagal baca html: %s" % str(e)[:120]
    if BeautifulSoup is None:
        teks = _clean(re.sub(r"<[^>]+>", " ", html))
        return "", seksi_dari_teks(teks), teks, ""
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    judul = ""
    try:
        if soup.title and soup.title.get_text().strip():
            judul = soup.title.get_text().strip()
    except Exception:
        judul = ""
    sections, all_txt = [], []
    cur_head, cur = "", []

    def flush():
        isi = _clean("\n".join(cur))
        if _n_nonspace(isi) >= MIN_ISI or (isi and cur_head):
            sections.append({"bagian": cur_head, "isi": isi})
            all_txt.append((cur_head + "\n" + isi).strip())

    body = soup.body or soup
    for el in body.find_all(["h1", "h2", "h3", "h4", "p", "li", "td", "pre"]):
        name = el.name.lower()
        txt = _clean(el.get_text(" "))
        if not txt:
            continue
        if name in ("h1", "h2", "h3", "h4"):
            flush()
            cur = []
            cur_head = txt[:90]
            if not judul:
                judul = txt
        else:
            cur.append(txt)
    flush()
    if not sections:
        teks = _clean(soup.get_text(" "))
        sections = seksi_dari_teks(teks)
        all_txt = [teks]
    return judul, sections, _clean("\n\n".join(all_txt)), ""


def _judul_dari_sections(sections, teks):
    for s in sections:
        b = _bersih_judul(s.get("bagian") or "")
        if judul_bermakna(b):
            return b
    for ln in (teks or "").splitlines():
        s = _bersih_judul(ln)
        if judul_bermakna(s):
            return s
    return ""


def classify(path: str, do_ocr: bool = False) -> SopFileInfo:
    """Deteksi format, ekstrak teks + bagian, tebak judul & kategori."""
    info = SopFileInfo(path=path)
    low = path.lower()
    base = os.path.basename(low)
    ext = os.path.splitext(low)[1]

    if base in SKIP_NAMES:
        info.tipe, info.catatan = "unknown", "berkas sistem (diabaikan)"
        return info
    if ext in LAWAS_EXT:
        info.tipe = "lawas"
        info.catatan = "format lama %s; konversi dulu ke .pptx/.docx" % ext
        return info

    judul_meta = ""
    try:
        if ext == ".pdf":
            info.tipe = "pdf"
            teks, perlu = _extract_pdf(path, do_ocr)
            info.perlu_ocr = perlu
            info.teks = _clean(teks)
            info.sections = seksi_dari_teks(info.teks)
            if perlu and not PF.has_tesseract():
                info.catatan = "PDF scan; OCR belum aktif (pasang Tesseract 'ind' + Poppler)"
        elif ext == ".pptx":
            info.tipe = "pptx"
            judul_meta, secs, teks, err = _extract_pptx(path)
            info.sections, info.teks, info.catatan = secs, teks, err
        elif ext == ".docx":
            info.tipe = "docx"
            judul_meta, secs, teks, err = _extract_docx(path)
            info.sections, info.teks, info.catatan = secs, teks, err
        elif ext in (".html", ".htm"):
            info.tipe = "html"
            judul_meta, secs, teks, err = _extract_html(path)
            info.sections, info.teks, info.catatan = secs, teks, err
        elif ext in (".txt", ".md"):
            info.tipe = "txt"
            info.teks = _clean(_extract_txt(path))
            info.sections = seksi_dari_teks(info.teks)
        elif ext in IMG_EXT:
            info.tipe = "image"
            teks, perlu = _extract_image(path, do_ocr)
            info.perlu_ocr = perlu
            info.teks = _clean(teks)
            info.sections = seksi_dari_teks(info.teks)
            if perlu and not PF.has_tesseract():
                info.catatan = "gambar; OCR belum aktif (butuh Tesseract 'ind')"
        else:
            info.tipe = "unknown"
            info.catatan = "format belum didukung"
            return info
    except Exception as e:  # pragma: no cover
        info.tipe = info.tipe or "unknown"
        info.catatan = ("gagal ekstrak: %s" % str(e))[:200]
        return info

    # Judul kandidat: metadata/heading dokumen -> heading pertama -> nama berkas.
    kand = _bersih_judul(judul_meta)
    if not judul_bermakna(kand):
        kand = _judul_dari_sections(info.sections, info.teks)
    if not judul_bermakna(kand):
        kand = _nama_dari_file(path)
    info.judul_kandidat = kand
    info.kategori_kandidat = tebak_kategori(os.path.basename(path), info.teks)
    return info
